import os
import argparse
from tqdm import tqdm
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract
import chromadb
from chromadb.utils import embedding_functions
import platform

def set_google_api_key():
    api_key = "AIzaSyD7VrRJrSa3W7u0syiZpWldChRCTiWLp-4"
    if platform.system() == "Windows":
        os.environ["GOOGLE_API_KEY"] = api_key
    else:  # Assuming macOS or Linux
        os.environ["GOOGLE_API_KEY"] = api_key
set_google_api_key()
def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text


def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text


def extract_text_from_image(image_path):
    """Extract text from an image using OCR."""
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text


def extract_text_from_file(file_path):
    """Extract text based on the file type."""
    _, extension = os.path.splitext(file_path)

    if extension.lower() == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension.lower() in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(file_path)
    elif extension.lower() == ".txt":
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read()
    else:
        print(f"Unsupported file type: {extension}")
        return None

def main(
    documents_directory: str = "./documents",
    collection_name: str = "documents_collection",
    persist_directory: str = ".",
) -> None:
    # Read all files in the documents directory
    documents = []
    metadatas = []
    files = os.listdir(documents_directory)
    for filename in files:
        file_path = os.path.join(documents_directory, filename)
        text = extract_text_from_file(file_path)

        if text:
            lines = text.splitlines()
            for line_number, line in enumerate(
                tqdm(lines, desc=f"Processing {filename}"), 1
            ):
                # Strip whitespace and skip empty lines
                line = line.strip()
                if len(line) == 0:
                    continue

                documents.append(line)
                metadatas.append({"filename": filename, "line_number": line_number})

    # Instantiate a persistent chroma client in the persist_directory.
    client = chromadb.PersistentClient(path=persist_directory)

    # Check for Google API Key or use Hugging Face embeddings
    """google_api_key = None
    if "GOOGLE_API_KEY" not in os.environ:
        gapikey = input("Please enter your Google API Key: ")
        genai.configure(api_key=gapikey)
        google_api_key = gapikey
    else:
        google_api_key = os.environ["GOOGLE_API_KEY"]"""

    # Example of using Hugging Face embedding function
    embedding_function = embedding_functions.HuggingFaceEmbeddingFunction(
        api_key="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",
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    # If the collection already exists, load it, or create a new one
    collection = client.get_or_create_collection(
        name=collection_name, embedding_function=embedding_function
    )

    # Create unique document IDs based on the current count
    count = collection.count()
    print(f"Collection already contains {count} documents")
    ids = [str(i) for i in range(count, count + len(documents))]

    # Load the documents in batches of 100
    for i in tqdm(
        range(0, len(documents), 100), desc="Adding documents", unit_scale=100
    ):
        collection.add(
            ids=ids[i : i + 100],
            documents=documents[i : i + 100],
            metadatas=metadatas[i : i + 100],  # type: ignore
        )

    new_count = collection.count()
    print(f"Added {new_count - count} documents")


if __name__ == "__main__":
    # Set up command line argument parsing
    parser = argparse.ArgumentParser(
        description="Load documents from a directory into a Chroma collection"
    )

    # Add arguments
    parser.add_argument(
        "--documents_directory",
        type=str,
        default="./documents",
        help="The directory where your files (PDFs, PPTs, images, text) are stored",
    )
    parser.add_argument(
        "--collection_name",
        type=str,
        default="documents_collection",
        help="The name of the Chroma collection",
    )
    parser.add_argument(
        "--persist_directory",
        type=str,
        default="chroma_storage",
        help="The directory where you want to store the Chroma collection",
    )

    # Parse arguments and run the main function
    args = parser.parse_args()

    main(
        documents_directory=args.documents_directory,
        collection_name=args.collection_name,
        persist_directory=args.persist_directory,
    )
