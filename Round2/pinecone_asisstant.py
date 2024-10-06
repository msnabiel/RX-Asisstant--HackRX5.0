import argparse
import os
from typing import List, Dict
from flask import Flask, request, jsonify
import requests
import google.generativeai as genai
from sentence_transformers import SentenceTransformer
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract
import platform
from flask_cors import CORS
from transformers import T5ForConditionalGeneration, T5Tokenizer
import warnings
from pinecone import Pinecone
import time
from pinecone_plugins.assistant.models.chat import Message

# Suppress specific warnings from Hugging Face transformers library
warnings.filterwarnings("ignore", message="You are using the default legacy behaviour")
warnings.filterwarnings("ignore", message="It will be set to `False` by default.")
warnings.filterwarnings("ignore", message="`clean_up_tokenization_spaces` was not set")

# Set Google API key for Gemini
def set_google_api_key():
    api_key = "AIzaSyD7VrRJrSa3W7u0syiZpWldChRCTiWLp-4"
    os.environ["GOOGLE_API_KEY"] = api_key

set_google_api_key()

# Initialize Flask app
app = Flask(__name__)
CORS(app)
os.environ["TOKENIZERS_PARALLELISM"] = "false"
# Retrieve the API key from the environment variable when needed
google_api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=google_api_key)

# Load the Flan-T5 model and tokenizer from Hugging Face
flan_tokenizer = T5Tokenizer.from_pretrained("google/flan-t5-base", legacy=False)
flan_model = T5ForConditionalGeneration.from_pretrained("google/flan-t5-base")

# Initialize history storage
session_history: Dict[str, List[Dict[str, str]]] = {}

# Load Huggingface model for embeddings
embedding_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

api_key = "7968bf7e-97f5-4022-adf7-9294590606be"

# Configure Pinecone client
pc = Pinecone(api_key=api_key)

# Create Assistant (if it doesn't already exist)
assistant_name = "example-assistant"
assistants = pc.assistant.list_assistants()
if assistant_name not in [a.name for a in assistants]:
    assistant = pc.assistant.create_assistant(
        assistant_name=assistant_name,
        instructions="You are a helpful and polite assistant.",
        timeout=30
    )
else:
    assistant = pc.assistant.Assistant(assistant_name=assistant_name)

# Helper functions to extract text from different document types
def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(image_path):
    """Extract text from an image using OCR."""
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_text_from_file(file_path):
    """Extract text based on the file type."""
    _, extension = os.path.splitext(file_path)
    if extension.lower() == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension.lower() in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(file_path)
    else:
        return None

def classify_query_with_flan(query: str, actions_list: list = ["create_order", "cancel_order", "collect_payment", "view_invoice"]) -> str:
    """Use Flan-T5 to classify the query as action-based or context-based."""
    
    # Dynamic prompt using the actions list
    actions_str = ', '.join(actions_list)
    prompt = f"""
    Classify the following query strictly as one of the actions: {actions_str}, or context-based.
    
    Query: {query}
    """
    print("Query for FlanT5:", query)
    inputs = flan_tokenizer(prompt, return_tensors="pt")
    outputs = flan_model.generate(**inputs, max_length=10, num_return_sequences=1)
    response = flan_tokenizer.decode(outputs[0], skip_special_tokens=True).lower()
    print("Response from FlanT5:", response)
    
    for action in actions_list:
        if action in response:
            return action
    return "context_based"

def build_combined_prompt(query: str, context: List[str], history: List[Dict[str, str]]) -> str:
    """Create a combined prompt to pass to the assistant based on context and history."""
    base_prompt = """
        Answer the user's question strictly based on the given context. If there is not enough information in the context, try your best to guess based on context.
        """
    user_prompt = f"The question is '{query}'. Here is all the context you have: {' '.join(context)}"
    history_prompt = "\n".join([f"User: {item['query']}\nBot: {item['response']}" for item in history])

    return f"{base_prompt} {history_prompt} {user_prompt}"

def get_assistant_response(query: str, context: List[str], session_id: str, document_id: str) -> str:
    """Get the assistant response, classify the query, and act accordingly."""
    history = session_history.get(session_id, [])

    # Classify the query using Flan-T5
    action = classify_query_with_flan(query)
    if action != "context_based":
        action_response = execute_action(action)
        session_history.setdefault(session_id, []).append({"query": query, "response": action_response})
        return action_response

    # Build the combined prompt
    prompt = build_combined_prompt(query, context, history)
    msg = Message(content=prompt)
    response = assistant.chat(messages=[msg])

    # Save the query and response in session history
    session_history.setdefault(session_id, []).append({"query": query, "response": response.text})

    # Add references to the response
    references = "\n".join([f"From document '{document_id}': Line {i + 1}: {line}" 
                            for i, line in enumerate(context)])
    return f"{response.text}\n\nReferences:\n{references}"

@app.route('/upload_document', methods=['POST'])
def upload_document():
    """Endpoint to upload and process a document."""
    uploads_dir = "uploads"
    os.makedirs(uploads_dir, exist_ok=True)

    if 'document' not in request.files:
        return jsonify({"error": "No document uploaded."}), 400

    document = request.files['document']
    doc_name = document.filename

    if not doc_name.endswith(('.pdf', '.ppt', '.pptx', '.jpg', '.jpeg', '.png', '.bmp', '.tiff')):
        return jsonify({"error": "Unsupported file format."}), 400

    file_path = os.path.join(uploads_dir, doc_name)
    document.save(file_path)

    extracted_text = extract_text_from_file(file_path)

    if extracted_text is None:
        return jsonify({"error": "Unable to extract text from the document."}), 400

    # Upload the document to Pinecone Assistant
    assistant.upload_file(file_path=file_path)

    return jsonify({"message": "Document uploaded and processed successfully."}), 200

@app.route('/chat', methods=['POST'])
def chat():
    """Endpoint to chat with the assistant."""
    data = request.get_json()
    query_text = data.get("query")
    session_id = data.get("session_id")
    document_id = data.get("document_id")

    if not query_text or not session_id:
        return jsonify({"error": "Missing query or session ID."}), 400

    query_vector = embedding_model.encode(query_text).tolist()

    if document_id:
        context_vectors = assistant.query(query_vector, top_k=5, include_metadata=True)
        context = [match["metadata"].get("line", "No text available") for match in context_vectors["matches"] if match["id"].startswith(document_id)]
    else:
        context_vectors = assistant.query(query_vector, top_k=5, include_metadata=True)
        context = [match["metadata"].get("line", "No text available") for match in context_vectors["matches"]]

    response = get_assistant_response(query_text, context, session_id, document_id)

    return jsonify({"response": response}), 200

if __name__ == '__main__':
    app.run(debug=True, port=5000)
