import argparse
import os
from typing import List, Dict
from flask import Flask, request, jsonify
import requests
import google.generativeai as genai
import chromadb
from chromadb.utils import embedding_functions
from sentence_transformers import SentenceTransformer
from tqdm import tqdm
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract
import platform
from flask_cors import CORS 
from transformers import T5ForConditionalGeneration, T5Tokenizer
import warnings
import json
from collections import defaultdict
conversation_state = defaultdict(dict)

# Suppress specific warnings from Hugging Face transformers library
warnings.filterwarnings("ignore", message="You are using the default legacy behaviour")
warnings.filterwarnings("ignore", message="It will be set to `False` by default.")
warnings.filterwarnings("ignore", message="`clean_up_tokenization_spaces` was not set")

# Set Google API key for Gemini
def set_google_api_key():
    api_key = "AIzaSyD7VrRJrSa3W7u0syiZpWldChRCTiWLp-4"
    os.environ["GOOGLE_API_KEY"] = api_key

set_google_api_key()

# Initialize Flask app
app = Flask(__name__)
CORS(app)
collection = None

# Now you can retrieve the API key from the environment variable when needed
google_api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=google_api_key)

UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
if not os.path.exists(UPLOAD_FOLDER):
    print("Creating folder...")
    os.makedirs(UPLOAD_FOLDER)
else:
    print("Folder already exists.")

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Load the Flan-T5 model and tokenizer from Hugging Face
flan_tokenizer = T5Tokenizer.from_pretrained("google/flan-t5-large", legacy=False)
flan_model = T5ForConditionalGeneration.from_pretrained("google/flan-t5-large")

# Initialize history storage
session_history: Dict[str, List[Dict[str, str]]] = {}

# Load Huggingface model for embeddings
embedding_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

# Initialize Gemini model
model = genai.GenerativeModel("gemini-1.5-flash")

# Helper functions to extract text from different document types
def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(image_path):
    """Extract text from an image using OCR."""
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_text_from_file(file_path):
    """Extract text based on the file type."""
    _, extension = os.path.splitext(file_path)
    if extension.lower() == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension.lower() in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(file_path)
    else:
        return None

def fetch_and_call_api(query):
    response_data = {}
    try:
        response = requests.get("https://dummyapi.com/api", params={"query": query})
        response_data = response.json()
    except Exception as e:
        return "Need API Key to call, to perform the action."

    if response and response_data.get("status") == "success":
        return response_data["message"]

def execute_action(action_name: str, query: str) -> str:
    """
    Executes the specified action based on the provided action name and query.
    """
    key_points_json = extract_key_action_with_gemini(query)
    extracted_data = json.loads(key_points_json)
    print("Extracted data:", extracted_data)

    # Filter out null values
    order_id = extracted_data.get("order_id")
    payment_id = extracted_data.get("payment_id")
    invoice_id = extracted_data.get("invoice_id")
    print("Order ID:", order_id, "Payment ID:", payment_id, "Invoice ID:", invoice_id)
    # Check for missing identifiers and prompt the user
    missing_info = []
    
    if action_name == "create_order" and order_id is None:
        missing_info.append("Order ID")
    if action_name == "cancel_order" and order_id is None:
        missing_info.append("Order ID")
    if action_name == "collect_payment" and payment_id is None:
        missing_info.append("Payment ID")
    if action_name == "view_invoice" and invoice_id is None:
        missing_info.append("Invoice ID")
    print(missing_info)
    if missing_info:
        return f"Please provide the following information: {', '.join(missing_info)}."

    # Proceed with the action if all required identifiers are present
    confirmation_message = ""
    if action_name == "create_order":
        return "Order created successfully."
    elif action_name == "cancel_order":
        confirmation_message = f"Are you sure you want to cancel your order with ID {order_id}?"
    elif action_name == "collect_payment":
        confirmation_message = f"Are you sure you want to collect the payment with ID {payment_id}?"
    elif action_name == "view_invoice":
        return f"Here is your invoice with ID {invoice_id}."
    else:
        return "No action taken."

    return confirmation_message

def confirm_action(action_name: str, identifier: str) -> str:
    """
    Confirms the specified action and executes it if confirmed.
    """
    if action_name == "cancel_order":
        return f"Your order with ID {identifier} has been cancelled."

    elif action_name == "collect_payment":
        return f"Payment with ID {identifier} has been collected successfully."

    # Additional actions can be confirmed here as needed
    return "No action taken."

# Define the default actions list outside the function
DEFAULT_ACTIONS_LIST = ["create_order", "cancel_order", "collect_payment", "view_invoice"]

def classify_query_with_flan(query: str, actions_list: list = DEFAULT_ACTIONS_LIST) -> str:
    """Use Flan-T5 to classify the query as action-based or context-based using a dynamic list of actions."""
    
    # Convert the list of actions to a string to use in the prompt
    actions_str = ', '.join(actions_list)
    
    # Dynamic prompt using the actions list
    prompt = f"""
    Classify the following query strictly as one of the actions: {actions_str}, or context-based.
    
    Query: {query}
    """
    
    # Tokenize the input
    inputs = flan_tokenizer(prompt, return_tensors="pt")
    outputs = flan_model.generate(**inputs, max_length=10, num_return_sequences=1)
    response = flan_tokenizer.decode(outputs[0], skip_special_tokens=True).lower()
    print(response)
    
    # Check if the response matches one of the actions
    for action in actions_list:
        if action in response:
            return action
    return "context_based"


def build_combined_prompt(query: str, context: List[str], history: List[Dict[str, str]]) -> str:
    base_prompt = """
        I am going to ask you a question, which I would like you to answer strictly based on the given context.
        If there is not enough information in the context or clear to answer the question, make a guess based on the context. YOUR RESPONSE SHOULD BE STRICTLY BASED ON CONTEXT.
        """
    user_prompt = f"The question is '{query}'. Here is all the context you have: {' '.join(context)}"
    history_prompt = "\n".join([f"User: {item['query']}\nBot: {item['response']}" for item in history])
    #print("Combined Prompt:", f"{base_prompt} {history_prompt} {user_prompt}")

    return f"{base_prompt} {history_prompt} {user_prompt}"

def classify_action_with_gemini(query: str) -> str:
    """
    Use Gemini to classify the query into an action or a context-based query.
    """
    prompt = f"""
    You are a helpful assistant. Classify the following query as one of the actions: create_order, cancel_order, collect_payment, view_invoice, 
    or context_based. 
    
    Query: "{query}"
    """
    
    # Send the prompt to Gemini for classification
    response = model.generate_content(prompt)
    classification = response.text.strip().lower()
    
    # Check if the classification matches one of the known actions
    for action in DEFAULT_ACTIONS_LIST:
        if action in classification:
            return action
    
    return "context_based"

def extract_key_action_with_gemini(query: str) -> str:
    """
    Use Gemini to extract key information from the query.
    """
    prompt = f"""
    You are a helpful assistant. Extract key information from the below statement, like order id, amount and return in json format. 
    
    statement: "{query}"
    """
    
    # Send the prompt to Gemini for classification
    response = model.generate_content(prompt)
    classification = response.text.strip().lower()
    #print("Gemini Key information before cleaning:", classification)
    classification = classification.replace("json", "").strip() # Remove "json" from the string and strip leading/trailing whitespace
    while classification.startswith('"') or classification.startswith("'") or classification.startswith("`"):
        classification = classification[1:]  # Remove the first character if it's a quote
    while classification.endswith('"') or classification.endswith("'") or classification.endswith("`"):
        classification = classification[:-1]  # Remove the last character if it's a quote

    # Remove backticks from the string
    classification = classification.replace("`", "")
    print("Gemini Key information after cleaning:", classification)
     # Convert string to dictionary to ensure valid JSON and remove any null values
    try:
        extracted_data = json.loads(classification)
        filtered_data = {k: v for k, v in extracted_data.items() if v is not None}
        return json.dumps(filtered_data)  # Return as JSON without null values
    except json.JSONDecodeError:
        print("JSON contains null values")
        return "{}"
    #return classification

def get_gemini_response(query: str, context: List[str], session_id: str, extract_data: bool = False) -> str:
    history = session_history.get(session_id, [])

    # Classify the query using Flan-T5
    action = classify_query_with_flan(query)
    action = action.lower()
    if action != "context_based":
        # If action is identified, extract necessary details
        action_response = execute_action(action, query)
        
        # Append the action execution to session history
        session_history.setdefault(session_id, []).append({"query": query, "response": action_response})
        return True, action_response

    # Build the combined prompt
    prompt = build_combined_prompt(query, context, history)
    print("SENT PROMPT TO GEMINI :", prompt)

    # Get response from Gemini
    response = model.generate_content(prompt)
    

    response_text = response.text.strip().lower()

    # Save the query and response in session history
    session_history.setdefault(session_id, []).append({"query": query, "response": response.text})

    return False, response.text


@app.route('/upload_document', methods=['POST'])
def upload_document():
    if 'document' not in request.files:
        return jsonify({"error": "No document uploaded."}), 400

    document = request.files['document']
    doc_name = document.filename

    # Validate document format
    if not doc_name.endswith(('.pdf', '.ppt', '.docx', '.jpg', '.jpeg', '.png')):
        return jsonify({"error": "Invalid document format. Only PDF, PPT, DOCX, and images are supported."}), 400

    try:
        # Save the uploaded file temporarily
        upload_path = os.path.join("uploads", doc_name)
        document.save(upload_path)

        # Extract text from the uploaded document
        text = extract_text_from_file(upload_path)
        if text is None:
            return jsonify({"error": f"Unsupported document format or failed to extract text from {doc_name}."}), 400

        # Generate a unique document ID
        document_id = "doc_" + os.urandom(6).hex()

        # Split text into lines for processing
        lines = text.splitlines()

        # Store extracted text into ChromaDB with metadata
        documents = [line for line in lines if line.strip()]
        #metadatas = [{"filename": doc_name, "line_number": i} for i, _ in enumerate(documents, 1)]
        metadatas = [{"document_id": document_id, "filename": doc_name, "line_number": i} for i, _ in enumerate(documents, 1)]

        # Add documents to Chroma collection
        collection.add(
            ids=[f"{document_id}_{i}" for i in range(len(documents))],
            documents=documents,
            metadatas=metadatas
        )

        return jsonify({"message": "Document uploaded and processed successfully.", "document_id": document_id}), 200

    except Exception as e:
        return jsonify({"error": f"An error occurred: {str(e)}"}), 500

@app.route('/chat', methods=['POST'])
def chat():
    data = request.json
    user_id = request.headers.get('x-user-id')
    session_id = request.headers.get('x-session-id')
    query = data.get('query')
    document_id = data.get('document_id')  # Optional document ID

    if not query:
        return jsonify({"error": "Query is required."}), 400

    try:
        if document_id:
            # Fetch context only from the specified document
            context_results = collection.query(
                query_texts=[query],
                n_results=5,
                include=["documents", "metadatas"],
                where={"document_id": document_id}  # Limit to specific document by ID
            )
        else:
            # Search across all documents when no document_id is provided
            context_results = collection.query(
                query_texts=[query],
                n_results=5,
                include=["documents", "metadatas"]  # No 'where' clause, consider the entire database
            )

        # If no relevant documents are found
        if not context_results["documents"]:
            return jsonify({"error": "No relevant information found in the database."}), 404

        # Extract the relevant context and metadata (line number and document name)
        # Flattening documents and metadata
        context = [doc for docs in context_results["documents"] for doc in docs]  # Flatten documents
        metadata = [meta for metas in context_results["metadatas"] for meta in metas]  # Flatten metadata

        # Check if context is a flat list of strings
        if not all(isinstance(c, str) for c in context):
            return jsonify({"error": "Context is not in the expected format."}), 500

        # Get response from Gemini model
        #response = get_gemini_response(query, context, session_id) + "\n\n"
        action_flag, response = get_gemini_response(query, context, session_id)

        if action_flag:
            return jsonify({"bot_message": response}), 200


        # Prepare the reference excerpts and metadata (line number and document name)
        references = []
        for i, (doc, meta) in enumerate(zip(context, metadata)):
            references.append({
                "excerpt": doc[:200],  # Show the first 200 characters of the relevant context
                "line_number": meta["line_number"],
                "document_name": meta["filename"]
            })

        # Format the response with relevant metadata
        response_data = {
            "bot_message": response,
            "references": references  # Include reference excerpts, line numbers, and document names
        }

        return jsonify(response_data), 200

    except Exception as e:
        return jsonify({"error": f"An internal error occurred: {str(e)}"}), 500


def main(collection_name: str = "documents_collection", persist_directory: str = ".") -> None:
    global collection  # Declare the global variable first

    client = chromadb.PersistentClient(path=persist_directory)

    # Create embedding function using Huggingface transformers
    embedding_function = embedding_functions.HuggingFaceEmbeddingFunction(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        api_key="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"
    )

    try:
        collection = client.get_collection(name=collection_name, embedding_function=embedding_function)
        print(f"Loaded existing collection: {collection_name}")
    except Exception as e:
        print(f"Collection '{collection_name}' not found. Creating a new collection.")
        collection = client.create_collection(name=collection_name, embedding_function=embedding_function)

    # Start the Flask app
    app.run(port=3000)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Load documents into a Chroma collection")
    parser.add_argument("--persist_directory", type=str, default="chroma_storage", help="Directory to store the Chroma collection")
    parser.add_argument("--collection_name", type=str, default="documents_collection", help="Name of the Chroma collection")

    args = parser.parse_args()

    main(collection_name=args.collection_name, persist_directory=args.persist_directory)
