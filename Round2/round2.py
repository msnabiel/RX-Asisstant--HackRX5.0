import argparse
import os
from typing import List, Dict
from flask import Flask, request, jsonify, g
import requests
import google.generativeai as genai
import chromadb
from chromadb.utils import embedding_functions
from sentence_transformers import SentenceTransformer
import os
import argparse
from tqdm import tqdm
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract
import platform
from flask_cors import CORS 

# Set Google API key for Gemini
def set_google_api_key():
    api_key = "AIzaSyD7VrRJrSa3W7u0syiZpWldChRCTiWLp-4"
    if platform.system() == "Windows":
        os.environ["GOOGLE_API_KEY"] = api_key
    else:  # Assuming macOS or Linux
        os.environ["GOOGLE_API_KEY"] = api_key

set_google_api_key()

# Initialize Flask app
app = Flask(__name__)
CORS(app)
collection = None
# Now you can retrieve the API key from the environment variable when needed
google_api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=google_api_key)

# Initialize history storage
session_history: Dict[str, List[Dict[str, str]]] = {}

# Load Huggingface model for embeddings
embedding_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

# Initialize Gemini model
model = genai.GenerativeModel("gemini-1.5-flash")

# Helper functions to extract text from different document types
def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(image_path):
    """Extract text from an image using OCR."""
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_text_from_file(file_path):
    """Extract text based on the file type."""
    _, extension = os.path.splitext(file_path)
    if extension.lower() == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension.lower() in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(file_path)
    else:
        return None


def build_combined_prompt(query: str, context: List[str], history: List[Dict[str, str]]) -> str:
    base_prompt = """
        I am going to ask you a question, which I would like you to answer strictly
        based on the given rules. 
        If the question or the intent involves any of the following actions, return the action name alone, you should not say anything else. Your response must be only the action name phrase:
        action names : create_order, cancel_order, collect_payment, view_invoice.
        Here are some examples:
        1. 'Create an order to book a bike' , say "create_order"
        2. 'Cancel my order', say "cancel_order"
        3. 'I want to pay for my purchase', say "collect_payment"
        4. 'Show me the invoice for order #123', say "view_invoice"
        If the question does not involve any of these actions, you need to answer strictly based on the given context.
        If there is not enough information in the context to answer the question, make a guess based on the context.
        """
    user_prompt = f"The question is '{query}'. Here is all the context you have: {' '.join(context)}"
    history_prompt = "\n".join([f"User: {item['query']}\nBot: {item['response']}" for item in history])

    return f"{base_prompt} {history_prompt} {user_prompt}"

def fetch_and_call_api(query):
    response_data = {}
    try:
        response = requests.get("https://dummyapi.com/api", params={"query": query})
        response_data = response.json()
    except Exception as e:
        return "Need API Key to call, to perform the action."

    if response and response_data.get("status") == "success":
        return response_data["message"]

def execute_action(action_name: str) -> str:
    if action_name == "create_order":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + "Order created successfully."
    elif action_name == "cancel_order":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + "Order cancelled successfully."
    elif action_name == "collect_payment":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + "Payment collected successfully."
    elif action_name == "view_invoice":
        fetch_response = fetch_and_call_api("YOUR_API_KEY")
        return fetch_response + "Here is your invoice."
    else:
        return "No action taken."

def get_gemini_response(query: str, context: List[str], session_id: str) -> str:
    history = session_history.get(session_id, [])

    # Build the combined prompt
    prompt = build_combined_prompt(query, context, history)

    # Get response from Gemini
    response = model.generate_content(prompt)

    response_text = response.text.strip().lower()

    # Check if Gemini returned an action instead of a regular answer
    if response_text in ["create_order", "cancel_order", "collect_payment", "view_invoice"]:
        action_response = execute_action(response_text)
        
        # Append the action execution to session history
        session_history.setdefault(session_id, []).append({"query": query, "response": action_response})
        return action_response

    # Save the query and response in session history
    session_history.setdefault(session_id, []).append({"query": query, "response": response.text})

    return response.text


@app.route('/upload_document', methods=['POST'])
def upload_document():
    collection = g.collection
    if 'document' not in request.files:
        return jsonify({"error": "No document uploaded."}), 400

    document = request.files['document']
    doc_name = document.filename

    # Validate document format
    if not doc_name.endswith(('.pdf', '.ppt', '.docx', '.jpg', '.jpeg', '.png')):
        return jsonify({"error": "Invalid document format. Only PDF, PPT, DOCX, and images are supported."}), 400

    try:
        # Save the uploaded file temporarily
        upload_path = os.path.join("uploads", doc_name)
        document.save(upload_path)

        # Extract text from the uploaded document
        text = extract_text_from_file(upload_path)
        if text is None:
            #app.logger.error(f"Failed to extract text from {doc_name}.")
            return jsonify({"error": "Unsupported document format or failed to extract text from {doc_name}."}), 400

        #if text is None:
        #    return jsonify({"error": "Unsupported document format or failed to extract text."}), 400

        # Generate a unique document ID
        document_id = "doc_" + os.urandom(6).hex()

        # Split text into lines for processing
        lines = text.splitlines()

        # Store extracted text into ChromaDB with metadata
        documents = [line for line in lines if line.strip()]
        metadatas = [{"filename": doc_name, "line_number": i} for i, _ in enumerate(documents, 1)]

        # Add documents to Chroma collection
        collection.add(
            ids=[f"{document_id}_{i}" for i in range(len(documents))],
            documents=documents,
            metadatas=metadatas
        )

        return jsonify({"document_id": document_id, "message": "Document uploaded and processed successfully."}), 200

    except Exception as e:
        return jsonify({"error": f"Failed to process document: {str(e)}"}), 500

    finally:
        # Clean up by removing the uploaded file after processing
        if os.path.exists(upload_path):
            os.remove(upload_path)


@app.route('/chat', methods=['POST'])
def chat():
    collection = g.collection
    data = request.json
    user_id = request.headers.get('x-user-id')
    session_id = request.headers.get('x-session-id')
    query = data.get('query')
    document_id = data.get('document_id')  # Optional document ID

    if not query:
        return jsonify({"error": "Query is required."}), 400

    try:
        if document_id:
            # Fetch context only from the specified document
            context_results = collection.query(
                query_texts=[query],
                n_results=5,
                include=["documents", "metadatas"],
                where={"document_id": document_id}  # Limit to specific document by ID
            )
        else:
            # Search across all documents when no document_id is provided
            context_results = collection.query(
                query_texts=[query],
                n_results=5,
                include=["documents", "metadatas"]  # No 'where' clause, consider the entire database
            )

        # If no relevant documents are found
        if not context_results["documents"]:
            return jsonify({"error": "No relevant information found in the database."}), 404

        # Extract the relevant context and metadata (line number and document name)
        context = [doc for sublist in context_results["documents"] for doc in sublist]
        metadata = [meta for sublist in context_results["metadatas"] for meta in sublist]

        # Get response from Gemini model
        response = get_gemini_response(query, context, session_id)

        # Prepare the reference excerpts and metadata (line number and document name)
        references = []
        for i, (doc, meta) in enumerate(zip(context, metadata)):
            references.append({
                "excerpt": doc[:200],  # Show the first 200 characters of the relevant context
                "line_number": meta["line_number"],
                "document_name": meta["filename"]
            })

        # Format the response with relevant metadata
        response_data = {
            "bot_message": response,
            "references": references  # Include reference excerpts, line numbers, and document names
        }

        return jsonify(response_data), 200

    except Exception as e:
        return jsonify({"error": "An internal error occurred. Please try again later."}), 500

def main(collection_name: str = "documents_collection", persist_directory: str = ".") -> None:
    global collection  # Declare the collection variable to access it in the chat function
    client = chromadb.PersistentClient(path=persist_directory)

    # Create embedding function using Huggingface transformers
    embedding_function = embedding_functions.HuggingFaceEmbeddingFunction(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        api_key="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"  # Ensure you set this environment variable
    )
    
    # Get the collection
    #collection = client.get_collection(name=collection_name, embedding_function=embedding_function)
    g.collection = client.get_collection(name=collection_name, embedding_function=embedding_function)

    # Start Flask app
    app.run(port=3000)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Load documents into a Chroma collection")
    parser.add_argument("--persist_directory", type=str, default="chroma_storage", help="Directory to store the Chroma collection")
    parser.add_argument("--collection_name", type=str, default="documents_collection", help="Name of the Chroma collection")

    args = parser.parse_args()

    main(collection_name=args.collection_name, persist_directory=args.persist_directory)